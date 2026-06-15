#!/usr/bin/env python3
"""Build the LiteLLM max_tokens investigation slide deck (PPTX)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0xC0, 0x4A, 0x2B)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF4, 0xF1, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x26, 0x2C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_box(slide, l, t, w, h, fill=None):
    box = slide.shapes.add_shape(1, l, t, w, h)
    box.line.fill.background()
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    box.shadow.inherit = False
    return box


def set_text(tf, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for spec in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = spec.get("align", align)
        p.space_after = Pt(spec.get("space_after", 6))
        p.space_before = Pt(spec.get("space_before", 0))
        if spec.get("bullet"):
            p.level = spec.get("level", 0)
        r = p.add_run()
        r.text = ("•  " if spec.get("bullet") else "") + spec["text"]
        f = r.font
        f.size = Pt(spec.get("size", 18))
        f.bold = spec.get("bold", False)
        f.italic = spec.get("italic", False)
        f.color.rgb = spec.get("color", DARK)
        f.name = "Calibri"


def content_slide(title, bullets, notes=None, kicker=None):
    s = prs.slides.add_slide(BLANK)
    add_box(s, 0, 0, SW, SH, LIGHT)
    add_box(s, 0, 0, Inches(0.28), SH, ACCENT)
    # title band
    tb = add_box(s, Inches(0.7), Inches(0.45), Inches(12.2), Inches(1.15))
    set_text(tb.text_frame,
             ([{"text": kicker, "size": 13, "bold": True, "color": ACCENT, "space_after": 2}] if kicker else [])
             + [{"text": title, "size": 30, "bold": True, "color": NAVY}],
             anchor=MSO_ANCHOR.MIDDLE)
    # body
    body = add_box(s, Inches(0.85), Inches(1.75), Inches(11.6), Inches(5.2))
    runs = []
    for b in bullets:
        if isinstance(b, tuple):
            text, lvl = b
        else:
            text, lvl = b, 0
        runs.append({"text": text, "bullet": True, "level": lvl,
                     "size": 18 if lvl == 0 else 16,
                     "color": DARK if lvl == 0 else GREY,
                     "space_after": 10 if lvl == 0 else 6})
    set_text(body.text_frame, runs)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


# ---- Slide 1: Title
s = prs.slides.add_slide(BLANK)
add_box(s, 0, 0, SW, SH, NAVY)
add_box(s, 0, Inches(4.95), SW, Inches(0.09), ACCENT)
t = add_box(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.6))
set_text(t.text_frame, [
    {"text": "Silent max_tokens Truncation in LiteLLM", "size": 40, "bold": True, "color": WHITE, "space_after": 10},
    {"text": "Root Cause Analysis and Remediation Plan", "size": 24, "color": RGBColor(0xC9, 0xD2, 0xE0)},
])
sub = add_box(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.4))
set_text(sub.text_frame, [
    {"text": "Engineering Investigation  |  LiteLLM 1.88.1", "size": 15, "bold": True, "color": ACCENT, "space_after": 4},
    {"text": "Premature response truncation affecting agent clients; diagnosis, fix, and rollout", "size": 14, "color": RGBColor(0xAF, 0xB8, 0xC8)},
])
s.notes_slide.notes_text_frame.text = (
    "Frame this as a closed investigation with a fix already implemented and tested in our fork, "
    "not an open question. The meeting goal is sign-off on the rollout path, not debate about the diagnosis.")

# ---- Slide 2
content_slide(
    "Executive Summary",
    ["Agent responses were cut off mid-output; some requests failed outright",
     "LiteLLM was silently rewriting the client's max_tokens before sending to the provider",
     "Triggered by a stale heuristic that misclassifies modern high-context models as legacy shared-window models",
     "Only active when modify_params: true is set, so the blast radius is precisely known",
     "Fixed with a one-condition gate; three deployment options range from zero-code config to an upstream patch"],
    notes="Management wants the shape of the problem and that it is solved. Most important line: this is config-gated, "
          "so we can describe the blast radius precisely.",
    kicker="READ THIS IF NOTHING ELSE")

# ---- Slide 3
content_slide(
    "What Users Actually Saw",
    ['Responses terminating early with finish_reason: "length" (OpenAI) or stop_reason: "max_tokens" (Anthropic)',
     "Worse the longer an agent session ran, because the cut scales with prompt size",
     "A second, more disruptive variant: hard failure with ContextWindowExceeded / provider 400",
     "Intermittent and config-dependent, which made it hard to pin down"],
    notes="Emphasize business impact: agents that appear to give up mid-task, broken tool calls, retries that cost "
          "tokens and latency. The two symptoms have different causes; we separate them later so we do not conflate them.",
    kicker="THE SYMPTOM")

# ---- Slide 4
content_slide(
    "The Log Signature That Confirmed It",
    ["With --detailed_debug, the proxy emits a literal MODIFYING MAX TOKENS line",
     "It prints the original user_max_tokens, the counted input_tokens, and the model's max_output_tokens",
     "Comparing the client-sent value against the raw payload leaving the proxy showed they differed",
     "This proved the truncation originated inside LiteLLM, not in the client or the provider"],
    notes="The smoking gun. Method: enable log_raw_request_response: true, capture one failing request, compare three "
          "numbers (client-sent max_tokens, payload-sent max_tokens, provider's reported figures). If client-sent equals "
          "payload-sent, LiteLLM is innocent for that request; if they differ, the rewrite fired.",
    kicker="THE PROOF")

# ---- Slide 5
content_slide(
    "Root Cause: get_modified_max_tokens in token_counter.py",
    ["When modify_params is on, LiteLLM recalculates max_tokens per request",
     'Case 1 fires only when model_info["max_input_tokens"] == max_output_tokens',
     "Assumption: equal limits mean a single shared window (like legacy gpt-3.5-turbo)",
     "It then subtracts the prompt: max_tokens = max_output - input_tokens - buffer",
     "The buffer is a thin max(0.1 * input_tokens, 10) based on a local token estimate"],
    notes="The equality check was a reasonable proxy in 2023 when shared-window completion models dominated. It encodes "
          "'input and output compete for the same budget.' The bug: this proxy is no longer reliable in the current registry.",
    kicker="ROOT CAUSE")

# ---- Slide 6
content_slide(
    "Why It Misfires on Modern Models",
    ["Hundreds of modern chat models list equal input/output limits but have independent budgets",
     "Examples include 256K/256K registry entries and our own 200,000-token configuration",
     "For these, LiteLLM wrongly subtracts the entire prompt from the output cap",
     "Result: max_tokens shrinks as the conversation grows, truncating long agent responses",
     "tiktoken undercounts non-OpenAI tokenizers by 15-20%, larger than the 10% safety buffer"],
    notes="Two failure modes. One, underflow risk: large prompts collapse the value toward zero (the code guards the "
          "strictly-negative case, but the value still ends up far below what the client asked for). Two, the buffer is "
          "provably smaller than tiktoken's typical undercount on Claude and Llama tokenizers, so even the margin is insufficient.",
    kicker="ROOT CAUSE")

# ---- Slide 7
content_slide(
    "Truncation vs Hard Failure: Different Causes",
    ["Silent truncation: directly caused by the LiteLLM rewrite shrinking max_tokens. Confirmed LiteLLM bug.",
     "Hard ContextWindowExceeded failure: the arithmetic is computed by the provider, with its real tokenizer, and relayed by LiteLLM",
     "LiteLLM does not contain a pre-flight prompt + max_tokens vs window off-by-one; the error is provider-enforced",
     "LiteLLM's real contribution to the hard case: its undercount-prone buffer can leave the rewritten value over budget on large prompts"],
    notes="This slide keeps us honest. The tempting but incorrect story is 'LiteLLM has an off-by-one that adds window plus "
          "one.' We audited every path; that off-by-one does not exist. The '+1 over' pattern is the fingerprint of a client "
          "computing fill-the-window max_tokens and overshooting by one token of estimation error, which the provider rejects. "
          "Present truncation as the LiteLLM bug we are fixing; present the hard failure as a client/provider boundary issue "
          "that our same fix also reduces. Do not cut this slide.",
    kicker="ACCURACY")

# ---- Slide 8
content_slide(
    "When Does This Actually Fire?",
    ["Only when litellm.modify_params (or LITELLM_MODIFY_PARAMS) is enabled",
     "Applies to chat completions, async completions, and the Anthropic /v1/messages path (Claude Code is in scope)",
     "Does not fire when modify_params is off, but turning it off globally disables useful behaviors",
     "Those behaviors include Anthropic tool-call repair and message sanitization",
     "Whether it does damage depends on the registry data for each deployment"],
    notes="Why a blanket 'set modify_params: false' is not the recommendation: the flag is overloaded, controlling both the "
          "buggy rewrite and a set of helpful Anthropic/Bedrock message-format fixes. We need a fix that disables only the rewrite.",
    kicker="BLAST RADIUS")

# ---- Slide 9
content_slide(
    "Option A: Correct the Model Registry",
    ["Override model_info per deployment so input and output limits are realistic and distinct",
     "Example: 200,000 input vs 4,096 or 64,000 output",
     "Breaking the equality means Case 1 can never trigger; only the harmless output cap remains",
     "Zero code change; pure YAML configuration",
     "Side benefit: improves routing decisions and cost tracking",
     ("Limitation: relies on every deployment being configured correctly; one stale entry reintroduces risk", 1)],
    notes="Lowest-risk, fastest mitigation; do this regardless of which deeper fix we adopt. Necessary but not sufficient "
          "alone, because it depends on configuration discipline across every model entry.",
    kicker="SOLUTION OPTION A")

# ---- Slide 10
content_slide(
    "Option B: Fix the Logic at the Source",
    ['Add one condition so the shrink only applies to genuine legacy models: mode == "completion" plus the equality check',
     "Verified against the live registry: legacy completion models still shrink; modern chat models with equal limits pass through",
     "Implemented in our fork with regression tests for both scenarios; lint and formatter clean",
     "Path to an upstream pull request so we stop carrying a private patch",
     "An environment-variable passthrough toggle was also considered as a more conservative alternative"],
    notes="Be transparent: we built the conditional mode gate, not the env-var toggle; the env var was a discussed fallback. "
          "Honest caveat for upstream reviewers: mode == 'completion' is a heuristic proxy, not a guarantee. If a truly "
          "shared-window model were ever registered as chat, it would skip the shrink and rely on provider enforcement, which "
          "is an acceptable trade because the provider enforces the window authoritatively anyway.",
    kicker="SOLUTION OPTION B")

# ---- Slide 11
content_slide(
    "Option C: Runtime Override in the Proxy",
    ["At proxy startup, replace the cached function pointer (_get_modified_max_tokens_func) with a passthrough",
     "Neutralizes only the rewrite; every other modify_params behavior is untouched",
     "Survives pip install -U litellm, unlike a source edit",
     "Zero change to the dependency itself; ideal as an immediate production stopgap",
     ("Limitation: an in-memory override is less visible to future engineers than a tracked code change", 1)],
    notes="Fastest way to protect production today without forking or waiting on an upstream merge. Clean, sanctioned patch "
          "point because the function is resolved through a lazy-import cache. Pair it with Option A so even models we forget "
          "to override are safe.",
    kicker="SOLUTION OPTION C")

# ---- Slide 12: matrix
s = prs.slides.add_slide(BLANK)
add_box(s, 0, 0, SW, SH, LIGHT)
add_box(s, 0, 0, Inches(0.28), SH, ACCENT)
tb = add_box(s, Inches(0.7), Inches(0.45), Inches(12.2), Inches(1.0))
set_text(tb.text_frame, [
    {"text": "DECISION MATRIX", "size": 13, "bold": True, "color": ACCENT, "space_after": 2},
    {"text": "Options Compared", "size": 30, "bold": True, "color": NAVY}], anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ["Option", "Effort", "Risk", "Scope", "Permanence"],
    ["A. Registry config", "Low", "Low", "Per-deployment", "Depends on discipline"],
    ["B. Upstream fix", "Medium", "Low", "Global, correct", "High; retires the fork"],
    ["C. Runtime patch", "Low", "Very low", "Global", "Immediate stopgap"],
]
left, top = Inches(0.85), Inches(1.85)
tbl_w, tbl_h = Inches(11.6), Inches(3.4)
table = s.shapes.add_table(len(rows), len(rows[0]), left, top, tbl_w, tbl_h).table
col_w = [Inches(3.0), Inches(1.7), Inches(1.8), Inches(2.6), Inches(2.5)]
for i, w in enumerate(col_w):
    table.columns[i].width = w
for r in range(len(rows)):
    for c in range(len(rows[0])):
        cell = table.cell(r, c)
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.08)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = rows[r][c]
        run.font.size = Pt(16 if r == 0 else 14)
        run.font.bold = (r == 0 or c == 0)
        if r == 0:
            run.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            run.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xEA, 0xE5, 0xDA)
note = add_box(s, Inches(0.85), Inches(5.5), Inches(11.6), Inches(1.3))
set_text(note.text_frame, [
    {"text": "These are complementary, not mutually exclusive: C protects production now, A hardens config and improves "
             "routing and cost data, B fixes it for everyone and lets us drop the fork.",
     "size": 15, "italic": True, "color": GREY}])
s.notes_slide.notes_text_frame.text = (
    "Frame the layering. We are not choosing one and discarding the others.")

# ---- Slide 13
content_slide(
    "Recommended Path: Layered Rollout",
    ["Immediately: deploy Option C (runtime passthrough) to stop production truncation today",
     "In parallel: apply Option A; correct model_info for our deployments, especially the 200,000-token config",
     "Near term: land Option B upstream so the fix is permanent and we retire the private patch",
     "Validate with a live proof: run the proxy with an equal input/output deployment, send a large max_tokens, confirm the raw provider payload carries the value unmodified"],
    notes="The validation step matters for sign-off. Per our conventions, proof-of-fix should be a real proxy hitting a "
          "real provider: before state logs MODIFYING MAX TOKENS with a shrunken value; after state forwards the client's "
          "value intact.",
    kicker="RECOMMENDATION")

# ---- Slide 14
content_slide(
    "Action Items",
    ["Apply the runtime patch to the production proxy config; confirm MODIFYING MAX TOKENS no longer appears",
     "Audit all model_info entries for input/output equality and correct them",
     "Open the upstream pull request from our fork branch with the conditional fix and regression tests",
     'Add monitoring that alerts if finish_reason: "length" rates rise again',
     "Decide whether to keep modify_params on (recommended, to retain Anthropic message-repair behaviors)"],
    notes="Assign an owner and a date to each line live in the meeting. The monitoring item turns this from a one-time fix "
          "into a regression guard.",
    kicker="NEXT STEPS")

# ---- Slide 15
content_slide(
    "Appendix: Technical Reference",
    ["File: litellm/litellm_core_utils/token_counter.py, function get_modified_max_tokens",
     "Wrapper gate: litellm/utils.py, the CHECK MAX TOKENS blocks, gated on modify_params and call type",
     "Patch point: litellm/_lazy_imports.py, _get_modified_max_tokens_func",
     "Version audited: 1.88.1; the fix and tests apply cleanly on our fork",
     "Diagnostic toggle: litellm_settings: log_raw_request_response: true"],
    notes="Leave this up during Q&A so any engineer can jump straight to the relevant code without interrupting the deck.",
    kicker="APPENDIX")

out = "/home/etty/LiteLLM_max_tokens_investigation.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
