#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build the Hebrew (RTL) LiteLLM max_tokens deck with implementation steps and the monkey patch."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0xC0, 0x4A, 0x2B)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF4, 0xF1, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x26, 0x2C)
CODEBG = RGBColor(0x16, 0x1B, 0x22)
CODEFG = RGBColor(0xD6, 0xDE, 0xE8)
CODEKEY = RGBColor(0xE6, 0x9A, 0x6B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_box(slide, l, t, w, h, fill=None):
    box = slide.shapes.add_shape(1, l, t, w, h)
    box.line.fill.background()
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    box.shadow.inherit = False
    return box


def _rtl(p, on=True):
    pPr = p._p.get_or_add_pPr()
    pPr.set("rtl", "1" if on else "0")


def set_text(tf, runs, anchor=MSO_ANCHOR.TOP):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for spec in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        rtl = spec.get("rtl", True)
        p.alignment = spec.get("align", PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT)
        _rtl(p, rtl)
        p.space_after = Pt(spec.get("space_after", 6))
        p.space_before = Pt(spec.get("space_before", 0))
        r = p.add_run()
        prefix = ("  •" if spec.get("bullet") and rtl else ("•  " if spec.get("bullet") else ""))
        # For RTL bullets put the marker on the right side of the text
        r.text = (spec["text"] + prefix) if (spec.get("bullet") and rtl) else (prefix + spec["text"])
        f = r.font
        f.size = Pt(spec.get("size", 18))
        f.bold = spec.get("bold", False)
        f.italic = spec.get("italic", False)
        f.color.rgb = spec.get("color", DARK)
        f.name = spec.get("font", "Arial")


def header(s, kicker, title):
    add_box(s, 0, 0, SW, SH, LIGHT)
    add_box(s, SW - Inches(0.28), 0, Inches(0.28), SH, ACCENT)  # rail on the right for RTL
    tb = add_box(s, Inches(0.7), Inches(0.45), Inches(11.9), Inches(1.15))
    set_text(tb.text_frame,
             [{"text": kicker, "size": 13, "bold": True, "color": ACCENT, "space_after": 2},
              {"text": title, "size": 30, "bold": True, "color": NAVY}],
             anchor=MSO_ANCHOR.MIDDLE)


def content_slide(kicker, title, bullets, notes=None):
    s = prs.slides.add_slide(BLANK)
    header(s, kicker, title)
    body = add_box(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(5.2))
    runs = []
    for b in bullets:
        text, lvl = b if isinstance(b, tuple) else (b, 0)
        runs.append({"text": text, "bullet": True,
                     "size": 19 if lvl == 0 else 16,
                     "color": DARK if lvl == 0 else GREY,
                     "space_after": 11 if lvl == 0 else 6})
    set_text(body.text_frame, runs)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def code_box(s, l, t, w, h, lines):
    add_box(s, l, t, w, h, CODEBG)
    inner = add_box(s, l + Inches(0.15), t + Inches(0.12), w - Inches(0.3), h - Inches(0.24))
    tf = inner.text_frame
    tf.word_wrap = True
    first = True
    for line, kind in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        _rtl(p, False)
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(13)
        r.font.color.rgb = CODEKEY if kind == "k" else CODEFG
    return inner


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK)
add_box(s, 0, 0, SW, SH, NAVY)
add_box(s, 0, Inches(4.95), SW, Inches(0.09), ACCENT)
t = add_box(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.7))
set_text(t.text_frame, [
    {"text": "קטיעת max_tokens שקטה ב-LiteLLM", "size": 40, "bold": True, "color": WHITE, "space_after": 10},
    {"text": "ניתוח שורש הבעיה, דרכי הפתרון ותוכנית ההטמעה", "size": 24, "color": RGBColor(0xC9, 0xD2, 0xE0)},
])
sub = add_box(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.4))
set_text(sub.text_frame, [
    {"text": "חקירה הנדסית  |  LiteLLM 1.88.1", "size": 15, "bold": True, "color": ACCENT, "space_after": 4},
    {"text": "קטיעת תשובות מוקדמת אצל סוכני AI: אבחון, תיקון והטמעה", "size": 14, "color": RGBColor(0xAF, 0xB8, 0xC8)},
])
s.notes_slide.notes_text_frame.text = (
    "להציג כחקירה סגורה עם תיקון שכבר מומש ונבדק ב-fork שלנו. מטרת הפגישה היא אישור מסלול ההטמעה, לא דיון באבחון.")

# ---------- Slide 2: Executive Summary ----------
content_slide("קראו את זה אם כלום אחר", "תקציר מנהלים", [
    "תשובות של סוכנים נקטעו באמצע; חלק מהבקשות נכשלו לחלוטין",
    "LiteLLM שכתב בשקט את max_tokens של הלקוח לפני השליחה לספק",
    "נגרם מהיוריסטיקה ישנה שמסווגת בטעות מודלים מודרניים עם חלון הקשר גדול כמודלי חלון משותף ישנים",
    "מופעל אך ורק כאשר מוגדר modify_params: true, ולכן רדיוס הפגיעה ידוע במדויק",
    "תוקן עם תנאי בודד; קיימות שלוש דרכי הטמעה, מהגדרה ללא קוד ועד תיקון ב-upstream",
], notes="ההנהלה רוצה את צורת הבעיה ושהיא נפתרה. השורה החשובה ביותר: הבאג מותנה בקונפיגורציה, ולכן אפשר לתאר את רדיוס הפגיעה במדויק.")

# ---------- Slide 3: Symptom ----------
content_slide("הסימפטום", "מה המשתמשים ראו בפועל", [
    'תשובות שנקטעו מוקדם עם finish_reason: "length" (פורמט OpenAI) או stop_reason: "max_tokens" (פורמט Anthropic)',
    "ככל שהשיחה ארוכה יותר, הקטיעה חמורה יותר, כי הקיצוץ גדל עם גודל הפרומפט",
    "וריאנט שני, מטריד יותר: כשל קשיח עם ContextWindowExceeded או 400 מהספק",
    "לא עקבי ותלוי-קונפיגורציה, ולכן היה קשה לאיתור",
], notes="להדגיש השפעה עסקית: סוכנים שנראה כאילו 'מוותרים' באמצע משימה, קריאות כלי שבורות, ניסיונות חוזרים שעולים בטוקנים ובזמן. שני הסימפטומים נובעים מסיבות שונות; נפריד ביניהם בהמשך.")

# ---------- Slide 4: Proof ----------
s = content_slide("ההוכחה", "חתימת הלוג שאישרה את הבאג", [
    "עם --detailed_debug, ה-proxy מדפיס שורה מילולית: MODIFYING MAX TOKENS",
    "השורה מציגה את user_max_tokens המקורי, את input_tokens שנספרו ואת max_output_tokens של המודל",
    "השוואה בין הערך שהלקוח שלח לבין ה-payload שיצא מה-proxy הראתה שהם שונים",
    "זה הוכיח שהקטיעה מקורה בתוך LiteLLM, לא בלקוח ולא בספק",
], notes="שקופית 'האקדח המעשן'. שיטה: להפעיל log_raw_request_response: true, ללכוד בקשה כושלת אחת, ולהשוות שלושה מספרים (max_tokens שהלקוח שלח, max_tokens שיצא ב-payload, והנתונים שהספק דיווח).")
code_box(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(1.5), [
    ("MODIFYING MAX TOKENS - user_max_tokens=200000, input_tokens=4210,", "n"),
    ("                       max_output_tokens=200000   ->  result=195790", "k"),
])

# ---------- Slide 5: Root cause ----------
s = content_slide("שורש הבעיה", "get_modified_max_tokens בקובץ token_counter.py", [
    "כש-modify_params פעיל, LiteLLM מחשב מחדש את max_tokens בכל בקשה",
    'Case 1 מופעל רק כאשר model_info["max_input_tokens"] == max_output_tokens',
    "ההנחה: גבולות שווים פירושם חלון משותף אחד (כמו gpt-3.5-turbo הישן)",
    "ואז הוא מחסיר את הפרומפט: max_tokens = max_output - input_tokens - buffer",
    "ה-buffer דק: max(0.1 * input_tokens, 10), מבוסס על הערכת טוקנים מקומית",
], notes="בדיקת השוויון הייתה פרוקסי סביר ב-2023 כשמודלי חלון משותף שלטו. הבאג: הפרוקסי הזה כבר לא אמין ברישום המודלים הנוכחי.")

# ---------- Slide 6: Why misfires ----------
content_slide("שורש הבעיה", "למה זה נכשל במודלים מודרניים", [
    "מאות מודלי chat מודרניים מציגים גבולות קלט ופלט שווים ברישום, אך בפועל יש להם תקציבים עצמאיים",
    "דוגמאות: רשומות 256K/256K, וכן הקונפיגורציה שלנו עם 200,000 טוקנים",
    "עבורם LiteLLM מחסיר בטעות את כל הפרומפט מתקרת הפלט",
    "התוצאה: max_tokens מתכווץ ככל שהשיחה גדלה, וקוטע תשובות ארוכות של סוכנים",
    "tiktoken מעריך בחסר מנגנוני טוקניזציה שאינם של OpenAI ב-15 עד 20 אחוז, יותר מה-buffer של 10 אחוז",
], notes="שני מצבי כשל. אחד, סיכון underflow: פרומפט גדול מקריס את הערך לכיוון אפס (הקוד שומר מפני שלילי ממש, אבל הערך עדיין יוצא הרבה מתחת למבוקש). שניים, ה-buffer מוכח כקטן מטעות החיסור של tiktoken על מנגנוני Claude ו-Llama.")

# ---------- Slide 7: Accuracy ----------
content_slide("דיוק", "קטיעה מול כשל קשיח: סיבות שונות", [
    "קטיעה שקטה: נגרמת ישירות מהשכתוב של LiteLLM שמכווץ את max_tokens. באג מאומת של LiteLLM.",
    "כשל קשיח ContextWindowExceeded: החשבון בהודעת השגיאה מחושב על-ידי הספק, עם הטוקנייזר האמיתי שלו, ורק מועבר דרך LiteLLM",
    "ב-LiteLLM אין בדיקת off-by-one מקדימה של prompt + max_tokens מול החלון; השגיאה נאכפת בצד הספק",
    "התרומה האמיתית של LiteLLM למקרה הקשיח: ה-buffer הנוטה להערכת חסר עלול להשאיר את הערך המשוכתב מעל התקציב בפרומפטים גדולים",
], notes="השקופית ששומרת עלינו מדויקים. הסיפור המפתה אך השגוי: 'ל-LiteLLM יש off-by-one שמוסיף חלון ועוד אחד'. בדקנו כל נתיב; ה-off-by-one הזה לא קיים. אל תוותרו על השקופית הזו.")

# ---------- Slide 8: Blast radius ----------
content_slide("רדיוס הפגיעה", "מתי זה באמת מופעל", [
    "רק כאשר litellm.modify_params (או LITELLM_MODIFY_PARAMS) מופעל",
    "חל על chat completions, async completions, ועל נתיב /v1/messages של Anthropic (כלומר Claude Code בתחום)",
    "לא מופעל כש-modify_params כבוי, אבל כיבוי גלובלי שלו משבית התנהגויות שימושיות אחרות",
    "ההתנהגויות האלה כוללות תיקון קריאות-כלי של Anthropic וניקוי הודעות",
    "האם זה גורם נזק תלוי בנתוני הרישום של כל deployment",
], notes="לכן 'modify_params: false' גורף אינו ההמלצה: הדגל עמוס מדי ושולט גם בשכתוב הבאגי וגם בתיקוני פורמט שימושיים של Anthropic/Bedrock. צריך תיקון שמשבית רק את השכתוב.")

# ---------- Slide 9: Option A with steps ----------
s = content_slide("פתרון א'", "תיקון רישום המודלים (Registry) ושלבי היישום", [
    "לדרוס model_info לכל deployment כך שגבולות הקלט והפלט יהיו ריאליים ושונים",
    "שלב 1: לאתר כל מודל שבו max_input_tokens שווה ל-max_output_tokens",
    "שלב 2: להגדיר ערכים נכונים, למשל 200,000 קלט מול 4,096 או 64,000 פלט",
    "שלב 3: לטעון מחדש את ה-proxy ולוודא ש-MODIFYING MAX TOKENS נעלם מהלוג",
    ("יתרון: אפס שינוי קוד, גם משפר ניתוב ומעקב עלויות. חיסרון: תלוי במשמעת קונפיגורציה בכל רשומה", 1),
], notes="המיטיגציה המהירה והבטוחה ביותר; לבצע בכל מקרה. הכרחי אך לא מספיק לבד.")
code_box(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.4), [
    ("model_list:", "k"),
    ("  - model_name: my-model", "n"),
    ("    litellm_params: { model: ... }", "n"),
    ("    model_info: { max_input_tokens: 200000, max_output_tokens: 4096 }", "k"),
])

# ---------- Slide 10: Option B with code ----------
s = content_slide("פתרון ב'", "תיקון הלוגיקה במקור (Upstream) ושלבי היישום", [
    'שלב 1: להוסיף תנאי שהכיווץ יחול רק על מודלים ישנים אמיתיים: mode == "completion" בנוסף לבדיקת השוויון',
    "שלב 2: אומת מול הרישום החי: מודלי completion ישנים עדיין מתכווצים נכון; מודלי chat עם גבולות שווים עוברים passthrough",
    "שלב 3: מומש ב-fork שלנו עם בדיקות רגרסיה לשני התרחישים; lint ו-formatter נקיים",
    "שלב 4: לפתוח Pull Request ל-upstream כדי שנפסיק להחזיק patch פרטי",
], notes="שקיפות: בנינו את תנאי ה-mode, לא toggle של משתנה סביבה (שנשקל כחלופה). הסתייגות הוגנת ל-reviewers: mode == 'completion' הוא פרוקסי היוריסטי, לא ערובה; אם מודל חלון-משותף אמיתי יירשם אי-פעם כ-chat, הוא ידלג על הכיווץ ויסתמך על אכיפת הספק, וזו עסקה מקובלת.")
code_box(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(1.5), [
    ("if (", "n"),
    ('    _model_info["max_input_tokens"] == max_output_tokens', "n"),
    ('    and _model_info.get("mode") == "completion"', "k"),
    ("):", "n"),
    ("    ...  # only legacy shared-window models shrink", "n"),
])

# ---------- Slide 11: Option C MONKEY PATCH ----------
s = prs.slides.add_slide(BLANK)
header(s, "פתרון ג'", "MONKEY PATCH בזמן ריצה (Runtime Override)")
body = add_box(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(2.0))
set_text(body.text_frame, [
    {"text": "מחליפים בעלייה את מצביע הפונקציה המטמון ל-passthrough; מנטרל רק את השכתוב, כל שאר התנהגויות modify_params נשמרות. שורד pip install -U litellm.", "bullet": True, "size": 18, "space_after": 8},
    {"text": "יישום: לשמור את הקובץ, ולהפנות אליו ב-litellm_settings: callbacks כך שייטען בעליית ה-proxy לפני הבקשה הראשונה.", "bullet": True, "size": 18, "color": GREY},
])
code_box(s, Inches(0.7), Inches(3.75), Inches(11.9), Inches(2.05), [
    ("# litellm_patch.py  -  loaded once at proxy startup", "n"),
    ("import litellm._lazy_imports as _li", "k"),
    ("", "n"),
    ("def _passthrough_max_tokens(*, user_max_tokens=None, **kwargs):", "k"),
    ("    return user_max_tokens", "n"),
    ("", "n"),
    ("_li._get_modified_max_tokens_func = _passthrough_max_tokens", "k"),
])
cfg = add_box(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.95))
set_text(cfg.text_frame, [
    {"text": "config.yaml:  litellm_settings -> callbacks: [\"litellm_patch\"]   (הקובץ חייב להיות ב-PYTHONPATH של ה-proxy)", "bullet": True, "size": 15, "color": NAVY, "rtl": False, "align": PP_ALIGN.LEFT},
])
s.notes_slide.notes_text_frame.text = (
    "הדרך המהירה ביותר להגן על הפרודקשן היום בלי fork ובלי להמתין למיזוג upstream. נקודת patch לגיטימית כי הפונקציה נפתרת דרך "
    "lazy-import cache. לשלב עם פתרון א' כדי שגם מודלים ששכחנו לדרוס יהיו מוגנים. הטיעון *kwargs קולט את כל ששת הארגומנטים שה-wrapper שולח.")

# ---------- Slide 12: Matrix ----------
s = prs.slides.add_slide(BLANK)
header(s, "מטריצת החלטה", "השוואת הפתרונות")
rows = [
    ["פתרון", "מאמץ", "סיכון", "היקף", "קביעות"],
    ["א. תיקון רישום", "נמוך", "נמוך", "לכל deployment", "תלוי במשמעת"],
    ["ב. תיקון upstream", "בינוני", "נמוך", "גלובלי ונכון", "גבוהה; מסיים את ה-fork"],
    ["ג. MONKEY PATCH", "נמוך", "נמוך מאוד", "גלובלי", "פתרון ביניים מיידי"],
]
left, top = Inches(0.7), Inches(1.9)
table = s.shapes.add_table(len(rows), 5, left, top, Inches(11.9), Inches(3.3)).table
for i, w in enumerate([Inches(3.0), Inches(1.7), Inches(1.9), Inches(2.7), Inches(2.6)]):
    table.columns[i].width = w
for r in range(len(rows)):
    for c in range(5):
        cell = table.cell(r, c)
        cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        _rtl(p, True)
        run = p.add_run(); run.text = rows[r][c]
        run.font.name = "Arial"
        run.font.size = Pt(16 if r == 0 else 14)
        run.font.bold = (r == 0 or c == 0)
        if r == 0:
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        else:
            run.font.color.rgb = DARK
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xEA, 0xE5, 0xDA)
note = add_box(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.2))
set_text(note.text_frame, [
    {"text": "הפתרונות משלימים, לא חלופיים: ג' מגן על הפרודקשן עכשיו, א' מחזק את הקונפיגורציה ומשפר ניתוב ועלויות, ב' מתקן לכולם ומאפשר לנו לוותר על ה-fork.",
     "size": 15, "italic": True, "color": GREY}])
s.notes_slide.notes_text_frame.text = "להציג את השכבתיות. איננו בוחרים אחד ופוסלים את האחרים."

# ---------- Slide 13: Recommendation ----------
content_slide("המלצה", "מסלול מומלץ: הטמעה שכבתית", [
    "מיידי: להטמיע את פתרון ג' (MONKEY PATCH) כדי לעצור את הקטיעה בפרודקשן היום",
    "במקביל: ליישם את פתרון א'; לתקן model_info ל-deployments שלנו, במיוחד לקונפיגורציה של 200,000 טוקנים",
    "בטווח הקרוב: למזג את פתרון ב' ל-upstream כדי שהתיקון יהיה קבוע ונסיר את ה-patch הפרטי",
    "לאמת בהוכחה חיה: להריץ proxy עם deployment של קלט=פלט, לשלוח max_tokens גדול, ולוודא שה-payload לספק נושא את הערך ללא שינוי",
], notes="שלב האימות חשוב לאישור. לפי המוסכמות שלנו, הוכחת תיקון צריכה להיות proxy אמיתי מול ספק אמיתי: מצב 'לפני' מתעד MODIFYING MAX TOKENS עם ערך מכווץ; מצב 'אחרי' מעביר את ערך הלקוח כמות שהוא.")

# ---------- Slide 14: Next steps ----------
content_slide("צעדים הבאים", "פריטי פעולה", [
    "להחיל את ה-MONKEY PATCH על קונפיגורציית הפרודקשן; לוודא ש-MODIFYING MAX TOKENS כבר לא מופיע",
    "לבדוק את כל רשומות model_info לאיתור שוויון קלט/פלט ולתקן אותן",
    "לפתוח Pull Request ל-upstream מענף ה-fork עם התיקון המותנה ובדיקות הרגרסיה",
    'להוסיף ניטור שמתריע אם שיעור finish_reason: "length" עולה שוב',
    "להחליט אם להשאיר את modify_params פעיל (מומלץ, כדי לשמר את תיקוני ההודעות של Anthropic)",
], notes="להקצות אחראי ותאריך לכל שורה בפגישה. פריט הניטור הופך את זה מתיקון חד-פעמי לשומר רגרסיה.")

# ---------- Slide 15: Appendix ----------
content_slide("נספח", "מקורות טכניים", [
    "קובץ: litellm/litellm_core_utils/token_counter.py, פונקציה get_modified_max_tokens",
    "שער ה-wrapper: litellm/utils.py, בלוקים של CHECK MAX TOKENS, מותנה ב-modify_params ובסוג הקריאה",
    "נקודת ה-patch: litellm/_lazy_imports.py, המשתנה _get_modified_max_tokens_func",
    "גרסה שנבדקה: 1.88.1; התיקון והבדיקות מוחלים נקי על ה-fork שלנו",
    "מתג אבחון: litellm_settings: log_raw_request_response: true",
], notes="להשאיר את השקופית פתוחה ב-Q&A כדי שכל מהנדס יוכל לקפוץ ישר לקוד הרלוונטי.")

out = "/home/etty/LiteLLM_max_tokens_investigation_HE.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
